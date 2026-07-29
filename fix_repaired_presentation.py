# -*- coding: utf-8 -*-
from pathlib import Path

from pptx import Presentation


SOURCE = Path(r"C:\Users\USER\Documents\secure-message-presentation.pptx [Repaired].pptx")
OUTPUT = SOURCE.with_name("secure-message-presentation-final-v3.pptx")


def replace_text(slide, old, new):
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False) and shape.text == old:
            paragraph = shape.text_frame.paragraphs[0]
            runs = list(paragraph.runs)
            if runs:
                # Assigning paragraph.text discards font, color, and run styling.
                # Reuse the original first run so the edited text keeps its design.
                runs[0].text = new
                for run in runs[1:]:
                    run._r.getparent().remove(run._r)
            else:
                paragraph.text = new
            return
    raise ValueError(f"Text not found on slide: {old!r}")


def main():
    presentation = Presentation(SOURCE)
    slides = presentation.slides

    # 1-3: keep the overview faithful to the current implementation.
    replace_text(slides[0], "허프만 트리 · AES-GCM · XOR 학습 모드 · 다익스트라 알고리즘", "허프만 트리 · AES-GCM · XOR 학습 모드 · 바이트 분포 분석")
    replace_text(slides[1], "허프만 / AES-GCM / XOR 학습 / 다익스트라", "허프만 / AES-GCM / XOR 학습 / 바이트 분포")
    replace_text(slides[2], "기밀성만이 아니라 용량, 무결성, 전달 경로까지 함께 고려해야 합니다.", "기밀성·무결성·용량을 함께 고려하고, 암호화 전후 데이터를 분석해야 합니다.")
    replace_text(slides[2], "경로", "분석")
    replace_text(slides[2], "같은 데이터라도 비용 기준에 따라\n좋은 전송 경로가 달라진다.", "암호화 전후의 바이트 분포와\nXOR 학습 결과를 비교한다.")
    replace_text(slides[2], "목표:  이산수학의 트리 · 부울대수 · 그래프를 실제 동작하는 웹앱에 연결한다.", "목표: 이산수학의 트리 · 부울대수 · 바이트 분포를 실제 웹앱에 연결한다.")

    # 2: replace the generic roadmap with the report's development purpose and real-world context.
    replace_text(slides[1], "ROADMAP", "PURPOSE & REAL-WORLD")
    replace_text(slides[1], "발표 흐름", "개발 목적과 실제 사례")
    replace_text(slides[1], "문제 정의부터 수학적 모델, 구현과 한계까지 설명합니다.", "이산수학 원리를 직접 구현하고, 실제 보안 기술과의 차이를 확인합니다.")
    replace_text(slides[1], "문제와 목표", "개발 목적")
    replace_text(slides[1], "왜 보안 전송을 수학적으로 모델링하는가", "압축·암호화·복호화 원리를 웹앱에서 직접 확인")
    replace_text(slides[1], "전체 구조", "실제 압축 사례")
    replace_text(slides[1], "입력부터 복호화·공유까지의 데이터 흐름", "ZIP·PNG의 DEFLATE는 허프만 코딩을 활용")
    replace_text(slides[1], "핵심 개념", "실제 암호 사례")
    replace_text(slides[1], "허프만 / AES-GCM / XOR 학습 / 바이트 분포", "HTTPS·파일 보호에는 AES-GCM 같은 인증 암호 사용")
    replace_text(slides[1], "구현과 시연", "교육용 구현")
    replace_text(slides[1], "웹앱 기능과 검증 과정", "브라우저에서 25MB 이하 파일의 처리 과정을 시각화")
    replace_text(slides[1], "한계와 확장", "실제 서비스와 차이")
    replace_text(slides[1], "공개 보관함의 한계 및 개선 방향", "로그인·권한·자동 삭제·악성 파일 검사 등이 추가로 필요")

    # 5: use a true variable-length, prefix-free Huffman code example.
    replace_text(slides[4], "허프만 코딩: 빈도가 코드 길이를 결정한다", "허프만 코딩: 빈도가 트리 구조를 결정한다")
    replace_text(slides[4], "예시: A = 00, B = 01, C = 10, D = 11  →  경계 표시 없이도 원문을 복원 가능", "표시된 트리 예시: A = 00, B = 01, C = 10, D = 11 → 경계 표시 없이 복원 가능")

    # 7: make the AAD header protection explicit, as described in the report.
    replace_text(slides[6], "복호화에 필요한 압축 정보, 키 유도 정보, 인증 암호문을 하나의 패키지로 전달합니다.", "복원 정보와 암호문을 하나로 전달하며, 헤더도 AAD로 인증해 변경을 거부합니다.")
    replace_text(slides[6], "변조·오류 시 즉시 거부", "암호문·헤더·키 오류 시 즉시 거부")

    # 8-9: replace nonexistent routing functionality with the implemented visuals and failure test.
    replace_text(slides[7], "09", "08")
    replace_text(slides[7], "동일 키 입력, 인증 검증 후 원본 복원", "허프만 트리와 압축 전후 크기 표시")
    replace_text(slides[7], "XOR 학습용 바이트 빈도 분포 비교", "암호화 전후·XOR 학습용 바이트 분포 비교")
    replace_text(slides[7], "복호화", "복호화·실패 검증")
    replace_text(slides[7], "전송 경로", "허프만 시각화")
    replace_text(slides[7], "시간/위험도 선택, 가중치 수정과 비교", "트리·코드표와 압축 전후 크기 표시")
    replace_text(slides[7], "Supabase에 암호문만 업로드·수신", "Supabase에 .enc 암호문만 업로드·수신")
    replace_text(slides[8], "경로 계산", "실패 검증")
    replace_text(slides[8], "시간 또는 위험도 기준으로 다익스트라 실행", "틀린 키 또는 변조 파일은 AES-GCM 인증 실패")
    replace_text(slides[8], "10", "09")

    # 10: report-derived operational limitations.
    replace_text(slides[9], "11", "10")

    # 11: retain Dijkstra only as a clearly-labelled future improvement.
    replace_text(slides[10], "03 / GRAPH", "08 / FUTURE IMPROVEMENT")
    replace_text(slides[10], "발전 방향 - 다익스트라: 가중치 합이 가장 작은 경로 찾기", "향후 개선 - 다익스트라로 전송 경로 최적화")
    replace_text(slides[10], "정점은 서버, 간선은 연결, 가중치는 전송 시간 또는 위험도로 모델링했습니다.", "현재 웹앱에는 미구현이며, 서버·연결·비용을 그래프로 모델링해 확장할 수 있습니다.")
    replace_text(slides[10], "시연", "향후 구현")
    replace_text(slides[10], "선택 경로와 탐색 기록, 패킷 이동을 표시.", "시간 또는 위험도 기준으로 최저 비용 경로를 탐색.")
    replace_text(slides[10], "08", "11")

    # 12: conclusion only claims implemented concepts; future routing is called out separately.
    replace_text(slides[11], "GRAPH", "XOR")
    replace_text(slides[11], "허프만 트리는 효율을, AES-GCM은 기밀성·무결성을, XOR은 부울대수 원리를, 다익스트라는 전달의 최적화를 설명합니다.", "허프만 트리는 효율을, AES-GCM은 기밀성·무결성을, XOR은 부울대수 원리를 설명합니다. 다익스트라는 향후 확장 아이디어입니다.")
    replace_text(slides[11], "“데이터의 표현·변환·이동은 모두 이산수학적 구조로 모델링할 수 있다.”", "“데이터의 표현·변환·보호는 이산수학적 구조로 모델링할 수 있다.”")

    presentation.save(OUTPUT)
    print(OUTPUT)


if __name__ == "__main__":
    main()
