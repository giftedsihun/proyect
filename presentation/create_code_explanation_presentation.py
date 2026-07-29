from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.dml.color import RGBColor


OUT = Path(__file__).resolve().parent / 'secure-message-code-explained-ko.pptx'
W, H = 13.333, 7.5
FONT = 'Malgun Gothic'
MONO = 'Consolas'
NAVY = '071924'
PANEL = '103442'
PANEL_DARK = '0A202C'
TEAL = '58E0CB'
GOLD = 'F2C85B'
CORAL = 'FF8977'
INK = 'EAF5F3'
MUTED = 'A6C3C7'
LINE = '2B5964'


def color(value):
    return RGBColor.from_string(value.lstrip('#'))


def text(slide, value, x, y, w, h, size=16, fill=INK, bold=False, font=FONT,
         align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = frame.margin_top = frame.margin_bottom = 0
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = value
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color(fill)
    return shape


def rect(slide, x, y, w, h, fill=PANEL, line=LINE, radius=True):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color(fill)
    shape.line.color.rgb = color(line)
    if radius:
        shape.adjustments[0] = 0.07
    return shape


def line(slide, x1, y1, x2, y2, fill=LINE, width=1.2):
    shape = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    shape.line.color.rgb = color(fill)
    shape.line.width = Pt(width)
    return shape


def circle(slide, x, y, d, fill, label='', label_fill=NAVY):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color(fill)
    shape.line.color.rgb = color(fill)
    if label:
        text(slide, label, x, y + d * .24, d, d * .4, 10, label_fill, True, MONO, PP_ALIGN.CENTER)
    return shape


def background(slide, section='SECURE MESSAGE / CODE EXPLAINED'):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color(NAVY)
    for x in [0.45 + i * .55 for i in range(25)]:
        line(slide, x, 0, x, H, '0B2B38', .25)
    for y in [0.35 + i * .55 for i in range(14)]:
        line(slide, 0, y, W, y, '0B2B38', .25)
    rect(slide, 0, 0, .16, H, TEAL, TEAL, False)
    text(slide, section, .48, .28, 6.8, .2, 7.5, MUTED, True, MONO)
    text(slide, 'CODE WALKTHROUGH', 10.1, .28, 2.7, .2, 7.5, MUTED, True, MONO, PP_ALIGN.RIGHT)


def footer(slide, number):
    line(slide, .62, 7.08, 12.7, 7.08, LINE, .7)
    text(slide, f'{number:02d}', .62, 7.18, .4, .2, 8, TEAL, True, MONO)
    text(slide, '실제 구현 기준: 허프만 트리 + PBKDF2 + AES-256-GCM', 1.02, 7.18, 5.6, .2, 8, MUTED)


def title(slide, kicker, heading, subtitle=''):
    text(slide, kicker, .62, .78, 6.5, .22, 9, TEAL, True, MONO)
    text(slide, heading, .6, 1.08, 12, .6, 28, INK, True)
    if subtitle:
        text(slide, subtitle, .62, 1.86, 11.5, .36, 13, MUTED)


def card(slide, x, y, w, h, kicker, heading, body, accent=TEAL):
    rect(slide, x, y, w, h)
    rect(slide, x + .18, y + .18, .06, h - .36, accent, accent, False)
    text(slide, kicker, x + .4, y + .27, w - .62, .18, 8, accent, True, MONO)
    text(slide, heading, x + .4, y + .57, w - .62, .3, 16, INK, True)
    text(slide, body, x + .4, y + 1.02, w - .62, h - 1.2, 11, MUTED)


def code_box(slide, x, y, w, h, code, accent=TEAL):
    rect(slide, x, y, w, h, PANEL_DARK, LINE)
    rect(slide, x + .18, y + .18, .05, h - .36, accent, accent, False)
    text(slide, code, x + .42, y + .25, w - .63, h - .45, 11.5, 'C9F3EB', False, MONO)


def add_flow(slide, items, y=3.0):
    width = 2.2
    for i, (head, body, accent) in enumerate(items):
        x = .62 + i * 2.52
        rect(slide, x, y, width, 1.55)
        circle(slide, x + .18, y + .18, .36, accent, f'{i + 1:02d}')
        text(slide, head, x + .22, y + .7, 1.75, .26, 13, INK, True)
        text(slide, body, x + .22, y + 1.05, 1.75, .3, 9.5, MUTED)
        if i < len(items) - 1:
            line(slide, x + width, y + .78, x + 2.46, y + .78, GOLD, 1.8)


def build():
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)
    blank = prs.slide_layouts[6]

    # 1. Cover
    s = prs.slides.add_slide(blank)
    background(s, 'SECURE MESSAGE / IMPLEMENTATION WALKTHROUGH')
    text(s, 'SECURE', .72, 1.15, 4.2, .55, 43, INK, True, MONO)
    text(s, 'MESSAGE', .72, 1.82, 5.2, .62, 46, TEAL, True, MONO)
    text(s, '코드로 이해하는 암호화 · 압축 · 복원', .76, 2.72, 8.5, .38, 23, INK, True)
    text(s, '이 프로그램이 실제로 어떤 순서로 데이터를 처리하는지 설명합니다.', .78, 3.3, 8.4, .28, 13, MUTED)
    items = [('입력', GOLD), ('허프만', TEAL), ('AES-GCM', CORAL), ('.enc', TEAL), ('복원', GOLD)]
    for i, (label, accent) in enumerate(items):
        x = 1.0 + i * 2.25
        circle(s, x, 4.45, .68, accent, str(i + 1))
        text(s, label, x - .2, 5.28, 1.1, .22, 11, INK, True, align=PP_ALIGN.CENTER)
        if i < len(items) - 1:
            line(s, x + .72, 4.79, x + 2.14, 4.79, TEAL, 2)
    rect(s, .76, 6.02, 8.25, .62, PANEL, LINE)
    text(s, '핵심: 원본과 비밀번호는 브라우저에서 처리하고, 공유 시 암호문만 보냅니다.', 1.02, 6.22, 7.75, .22, 12, INK, True)
    footer(s, 1)

    # 2. Scope
    s = prs.slides.add_slide(blank)
    background(s)
    title(s, '01 / PROJECT MAP', '이 프로그램은 무엇을 하나요?', '정적 HTML/CSS/JavaScript 웹앱이며 서버 없이도 암호화와 복원이 동작합니다.')
    card(s, .75, 2.6, 3.8, 2.55, 'INPUT', '텍스트 또는 파일', '텍스트 또는 25MB 이하 파일을 읽어 Uint8Array 바이트 데이터로 변환합니다.', GOLD)
    card(s, 4.78, 2.6, 3.8, 2.55, 'PROCESS', '압축 후 암호화', '허프만 압축 → 키 유도 → AES-GCM 인증 암호화를 순서대로 적용합니다.', TEAL)
    card(s, 8.81, 2.6, 3.8, 2.55, 'OUTPUT', '.enc 패키지', '복원에 필요한 메타데이터와 암호문을 SME2 형식 하나로 묶습니다.', CORAL)
    code_box(s, 1.25, 5.65, 10.8, .57, '원본 입력  ->  Huffman  ->  PBKDF2  ->  AES-256-GCM  ->  SME2 .enc', GOLD)
    footer(s, 2)

    # 3. File map
    s = prs.slides.add_slide(blank)
    background(s)
    title(s, '02 / FILE MAP', '어느 파일이 어떤 역할을 하나요?', '발표 중 코드 위치를 말할 때 이 구조를 기준으로 설명하면 됩니다.')
    card(s, .75, 2.45, 3.8, 1.55, 'UI', 'index.html', '암호화·복호화·분석·공유 보관함의 화면 구조와 입력 요소를 정의합니다.', TEAL)
    card(s, 4.78, 2.45, 3.8, 1.55, 'LOGIC', 'script.js', '허프만, AES-GCM, .enc 패키지, 그래프, Supabase 통신을 구현합니다.', GOLD)
    card(s, 8.81, 2.45, 3.8, 1.55, 'STYLE', 'style.css', '반응형 보안 대시보드 형태의 화면 디자인을 담당합니다.', CORAL)
    card(s, 2.77, 4.45, 3.8, 1.55, 'DOCUMENT', 'PROJECT_OVERVIEW_KO.md', '기능, 이산수학 연결, 보안 한계와 개발 계획을 설명합니다.', TEAL)
    card(s, 6.8, 4.45, 3.8, 1.55, 'PRESENTATION', 'presentation/', '기존 발표 자료, 포스터, 그리고 생성 Python 스크립트를 담고 있습니다.', GOLD)
    footer(s, 3)

    # 4. Encryption flow
    s = prs.slides.add_slide(blank)
    background(s)
    title(s, '03 / ENCRYPT FLOW', '암호화 버튼을 누르면 일어나는 일', '`encrypt()` 함수가 화면 상태와 데이터 처리 전체를 순서대로 제어합니다.')
    add_flow(s, [
        ('입력 읽기', '텍스트/파일 → bytes', GOLD),
        ('빈도 계산', '0~255 등장 횟수', TEAL),
        ('허프만 압축', '비트열로 축소', TEAL),
        ('키 유도', 'PBKDF2 + salt', CORAL),
        ('인증 암호화', 'AES-GCM + IV', CORAL),
    ])
    code_box(s, 1.0, 5.22, 11.25, .84, 'encrypt()\n  sourceFromInput() -> buildHuffman() -> compressHuffman() -> deriveAesKey() -> crypto.subtle.encrypt()', TEAL)
    footer(s, 4)

    # 5. Frequency and Huffman
    s = prs.slides.add_slide(blank)
    background(s)
    title(s, '04 / DISCRETE MATH', '빈도표와 허프만 트리', '가장 직접적으로 구현된 이산수학: 도수표, 최소 힙, 이진 트리, 접두어 없는 코드')
    code_box(s, .78, 2.55, 5.7, 1.4, 'frequencyTable(bytes)\n  frequencies[byte]++\n\nbuildHuffman(frequencies)\n  최소 빈도 노드 2개를 반복 결합', TEAL)
    rect(s, 7.1, 2.48, 5.1, 2.15)
    text(s, '작은 허프만 트리 예시', 7.42, 2.76, 3.7, .22, 11, TEAL, True, MONO)
    line(s, 9.62, 3.22, 8.28, 4.08, TEAL, 1.5)
    line(s, 9.62, 3.22, 10.96, 4.08, TEAL, 1.5)
    circle(s, 9.37, 2.96, .5, TEAL, '10')
    circle(s, 8.03, 3.84, .5, GOLD, 'A')
    circle(s, 10.71, 3.84, .5, GOLD, 'B')
    text(s, '0', 8.78, 3.46, .2, .2, 9, GOLD, True, MONO)
    text(s, '1', 10.25, 3.46, .2, .2, 9, GOLD, True, MONO)
    text(s, '왼쪽 0 / 오른쪽 1을 따라가며 코드 생성', 7.42, 4.34, 4.3, .2, 10, MUTED)
    card(s, .8, 4.92, 3.55, 1.22, 'WHY', '짧은 코드 배정', '자주 등장한 바이트일수록 더 짧은 비트 코드를 받습니다.', GOLD)
    card(s, 4.88, 4.92, 3.55, 1.22, 'PROPERTY', 'Prefix-free', '어떤 코드도 다른 코드의 앞부분이 아니어서 경계를 구분할 수 있습니다.', TEAL)
    card(s, 8.96, 4.92, 3.55, 1.22, 'RESTORE', '동일 트리 재생성', '패키지에 저장한 빈도표로 복호화 시 같은 트리를 다시 만듭니다.', CORAL)
    footer(s, 5)

    # 6. Compression and packaging
    s = prs.slides.add_slide(blank)
    background(s)
    title(s, '05 / BINARY DATA', '압축 데이터는 어떻게 .enc에 담기나요?', '문자열이 아니라 바이트 배열과 DataView를 이용해 순서가 정해진 이진 패키지를 만듭니다.')
    items = [('SME2', '형식 식별자', TEAL), ('이름/MIME', '원본 정보', GOLD), ('빈도표', '트리 재구성', TEAL), ('salt/IV', '키·암호화 값', CORAL), ('암호문', '태그 포함', GOLD)]
    for i, (head, body, accent) in enumerate(items):
        x = .65 + i * 2.52
        rect(s, x, 2.85, 2.1, 1.25)
        text(s, head, x + .18, 3.12, 1.7, .22, 13, accent, True, MONO)
        text(s, body, x + .18, 3.53, 1.7, .2, 10, MUTED)
    code_box(s, 1.02, 4.78, 11.18, .92, 'makeAesHeader(...)  ->  [SME2 | metadata | frequency table | salt | IV]\nmakeAesPackage(header, encrypted)  ->  header + AES-GCM ciphertext', GOLD)
    text(s, '복호화 시 parseAesPackage()가 같은 순서로 헤더를 읽습니다.', 1.05, 5.98, 9.2, .22, 11, MUTED)
    footer(s, 6)

    # 7. AES
    s = prs.slides.add_slide(blank)
    background(s)
    title(s, '06 / SECURITY', '기본 보안: PBKDF2 + AES-256-GCM', 'XOR은 비교 학습용이며, 실제 기본 모드는 인증 암호화인 AES-GCM입니다.')
    add_flow(s, [
        ('비밀번호', '사용자 입력', GOLD),
        ('새 salt', '매번 16 bytes', TEAL),
        ('PBKDF2', 'SHA-256 250,000회', CORAL),
        ('새 IV', '매번 12 bytes', TEAL),
        ('AES-GCM', '암호화 + 인증', CORAL),
    ], 2.7)
    code_box(s, .95, 4.82, 11.38, .86, 'deriveAesKey(password, salt)\ncrypto.subtle.encrypt({ name: "AES-GCM", iv, additionalData: header }, key, compressedBytes)', CORAL)
    rect(s, 1.4, 5.96, 10.5, .47, '162F37', '2B5964')
    text(s, 'header를 AAD로 함께 인증: 암호문뿐 아니라 파일 이름·빈도표 등 헤더 변조도 감지합니다.', 1.68, 6.1, 10, .2, 10.5, INK, True)
    footer(s, 7)

    # 8. Decrypt
    s = prs.slides.add_slide(blank)
    background(s)
    title(s, '07 / DECRYPT FLOW', '복호화는 왜 안전하게 실패하나요?', 'AES-GCM 인증 검증을 먼저 통과해야만 압축 해제와 원본 복원이 진행됩니다.')
    add_flow(s, [
        ('.enc 읽기', 'parsePackage()', GOLD),
        ('헤더 해석', 'salt/IV/빈도표', TEAL),
        ('키 재생성', '같은 PBKDF2', CORAL),
        ('GCM 검증', '실패 시 즉시 중단', CORAL),
        ('압축 해제', '원본 bytes 복원', TEAL),
    ], 2.8)
    code_box(s, .95, 4.96, 11.42, .9, 'crypto.subtle.decrypt({ name: "AES-GCM", iv: pack.iv, additionalData: pack.aad }, key, pack.encrypted)\n-> decompressHuffman(compressed, pack.bitLength, buildHuffman(pack.frequencies), pack.originalSize)', TEAL)
    text(s, '잘못된 비밀번호 또는 바뀐 파일/헤더 -> "AES-GCM 인증 검증에 실패" -> 복원하지 않음', 1.1, 6.18, 10.8, .22, 11.5, CORAL, True)
    footer(s, 8)

    # 9. XOR and visualizations
    s = prs.slides.add_slide(blank)
    background(s)
    title(s, '08 / LEARNING MODE', 'XOR과 바이트 분포 분석은 무엇을 보여주나요?', 'XOR은 부울대수 원리 비교용, 그래프는 바이트 값 0~255의 빈도 차이 관찰용입니다.')
    code_box(s, .75, 2.5, 5.6, 1.15, 'xorBytes(data, key)\n  result[i] = data[i] ^ key[i % key.length]\n\nA XOR K XOR K = A', CORAL)
    card(s, .75, 4.2, 3.65, 1.55, 'XOR', '자기역원 성질', '같은 키를 두 번 XOR하면 원래 바이트로 돌아오는 부울대수 성질을 확인합니다.', CORAL)
    card(s, 4.85, 4.2, 3.65, 1.55, 'GRAPH', '바이트 빈도 분포', 'frequencyTable() 결과를 Canvas 막대그래프로 그려 원본과 암호문을 비교합니다.', TEAL)
    card(s, 8.95, 4.2, 3.65, 1.55, 'WARNING', '실제 보호용 아님', '반복 XOR은 취약합니다. 새 파일 생성의 기본값은 AES-GCM입니다.', GOLD)
    text(s, '관련 함수: xorBytes(), drawDistribution(), renderAnalysis()', .9, 6.22, 9.6, .2, 11, MUTED, False, MONO)
    footer(s, 9)

    # 10. Share
    s = prs.slides.add_slide(blank)
    background(s)
    title(s, '09 / SHARING', 'Supabase 공유는 무엇을 보내고 무엇을 보내지 않나요?', '공유는 선택 기능이며, 브라우저가 암호화를 끝낸 후 .enc 패키지만 업로드합니다.')
    rect(s, .85, 2.6, 5.45, 2.5, '123842', '2B5964')
    text(s, '브라우저 안에 남는 것', 1.18, 2.95, 3.3, .25, 13, TEAL, True)
    text(s, '원본 파일\n비밀번호\n압축 전·후 데이터\nAES 키', 1.18, 3.44, 3.5, 1.25, 16, INK, True)
    rect(s, 7.03, 2.6, 5.45, 2.5, '332D22', '2B5964')
    text(s, 'Supabase에 올라가는 것', 7.36, 2.95, 3.6, .25, 13, GOLD, True)
    text(s, '.enc 암호문 패키지\n표시 이름\n패키지 크기\n만료 시각', 7.36, 3.44, 3.6, 1.25, 16, INK, True)
    line(s, 6.32, 3.85, 7.0, 3.85, TEAL, 2.5)
    text(s, '업로드', 6.27, 3.48, .8, .2, 9, TEAL, True, MONO, PP_ALIGN.CENTER)
    code_box(s, 1.1, 5.6, 11.1, .55, 'uploadEncryptedPackage() -> /storage/v1/object/...   +   /rest/v1/secure_messages', GOLD)
    footer(s, 10)

    # 11. Accurate claims
    s = prs.slides.add_slide(blank)
    background(s)
    title(s, '10 / PRESENTATION GUIDE', '발표에서 정확하게 말해야 할 것', '코드에 실제로 있는 기능과, 단지 사례로만 남은 내용을 구분하는 것이 중요합니다.')
    card(s, .75, 2.6, 5.75, 2.75, 'IMPLEMENTED', '실제로 구현됨', '허프만 이진 트리와 최소 힙\n바이트 빈도표와 Canvas 그래프\nPBKDF2 + AES-256-GCM\nSME2 .enc 패키지와 인증 복호화\nSupabase 암호문 공유', TEAL)
    card(s, 6.83, 2.6, 5.75, 2.75, 'NOT IN CURRENT APP', '실제 앱 기능은 아님', '다익스트라 최단 경로 계산\n서버 간 패킷 전송 시뮬레이션\n네트워크 경로 제어\n\n발표 자료의 현실 사례로만 설명해야 합니다.', CORAL)
    rect(s, 1.2, 5.9, 10.9, .47, '162F37', '2B5964')
    text(s, '한 문장: "이 앱은 트리·빈도표·XOR을 학습에 연결하고, AES-GCM으로 실제 암호화를 구현했습니다."', 1.47, 6.04, 10.35, .2, 10.5, INK, True)
    footer(s, 11)

    # 12. Closing
    s = prs.slides.add_slide(blank)
    background(s, 'CONCLUSION / SECURE MESSAGE')
    text(s, 'CODE', .76, 1.12, 2.1, .45, 29, TEAL, True, MONO)
    text(s, '+', 2.9, 1.1, .4, .45, 29, GOLD, True, MONO)
    text(s, 'DISCRETE MATH', 3.4, 1.12, 4.3, .45, 29, GOLD, True, MONO)
    text(s, '+', 7.75, 1.1, .4, .45, 29, CORAL, True, MONO)
    text(s, 'SECURITY', 8.3, 1.12, 3.2, .45, 29, CORAL, True, MONO)
    text(s, '코드의 흐름을 이해하면\n이산수학이 기능이 되는 과정이 보입니다.', .76, 2.0, 10.8, 1.2, 35, INK, True)
    rect(s, .8, 4.15, 11.7, 1.35)
    text(s, '핵심 정리', 1.08, 4.47, 1.3, .2, 10, TEAL, True, MONO)
    text(s, '빈도표와 트리는 압축을, PBKDF2와 AES-GCM은 보호를, .enc 패키지는 안전한 전달을 담당합니다.', 2.5, 4.4, 9.25, .35, 16, INK, True)
    text(s, '감사합니다', .8, 6.25, 3, .3, 18, TEAL, True, MONO)
    footer(s, 12)

    prs.save(OUT)
    print(OUT)


if __name__ == '__main__':
    build()
