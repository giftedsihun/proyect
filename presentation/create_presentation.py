from pathlib import Path
from math import cos, sin, pi
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR

ROOT = Path(__file__).resolve().parent
OUT = ROOT
W, H = 13.333, 7.5
FONT = 'Malgun Gothic'
MONO = 'Consolas'
NAVY = '071924'
NAVY2 = '0C2734'
PANEL = '103442'
TEAL = '58E0CB'
GOLD = 'F2C85B'
CORAL = 'FF8977'
INK = 'EAF5F3'
MUTED = 'A6C3C7'
LINE = '2B5964'

def rgb(value):
    value = value.lstrip('#')
    return RGBColor.from_string(value)

def add_text(slide, text, x, y, w, h, size=18, color=INK, bold=False, font=FONT, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear(); frame.word_wrap = True; frame.margin_left = 0; frame.margin_right = 0; frame.margin_top = 0; frame.margin_bottom = 0
    frame.vertical_anchor = valign
    p = frame.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = text
    run.font.name = font; run.font.size = Pt(size); run.font.bold = bold; run.font.color.rgb = rgb(color)
    return box

def rect(slide, x, y, w, h, fill=PANEL, line=None, radius=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(line or fill)
    if radius: shape.adjustments[0] = 0.08
    return shape

def line(slide, x1, y1, x2, y2, color=LINE, width=1.2, dash=None):
    shape = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    shape.line.color.rgb = rgb(color); shape.line.width = Pt(width)
    if dash: shape.line.dash_style = dash
    return shape

def circle(slide, x, y, d, fill=TEAL, label=None, label_color=NAVY, size=14):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    s.fill.solid(); s.fill.fore_color.rgb = rgb(fill); s.line.color.rgb = rgb(fill)
    if label:
        add_text(slide, label, x, y + d * .22, d, d * .45, size, label_color, True, align=PP_ALIGN.CENTER)
    return s

def bg(slide, section='SECURE MESSAGE / DISCRETE MATHEMATICS'):
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = rgb(NAVY)
    # faint asymmetric grid
    for x in [0.45 + i * .55 for i in range(25)]: line(slide, x, 0, x, H, '0B2B38', .25)
    for y in [0.35 + i * .55 for i in range(14)]: line(slide, 0, y, W, y, '0B2B38', .25)
    rect(slide, 0, 0, .16, H, TEAL)
    add_text(slide, section, .45, .28, 6.4, .25, 7.5, MUTED, True, MONO)
    add_text(slide, 'SECURE MESSAGE', 10.55, .28, 2.3, .25, 7.5, MUTED, True, MONO, PP_ALIGN.RIGHT)

def title(slide, kicker, heading, sub=''):
    add_text(slide, kicker, .62, .78, 5.3, .28, 9, TEAL, True, MONO)
    add_text(slide, heading, .6, 1.08, 11.9, .7, 29, INK, True)
    if sub: add_text(slide, sub, .62, 1.88, 11.5, .42, 13, MUTED)

def footer(slide, number):
    line(slide, .62, 7.08, 12.7, 7.08, LINE, .7)
    add_text(slide, f'{number:02d}', .62, 7.18, .4, .2, 8, TEAL, True, MONO)
    add_text(slide, '원본과 비밀 키는 브라우저 밖으로 나가지 않습니다.', 1.02, 7.18, 5, .2, 8, MUTED)

def card(slide, x, y, w, h, label, heading, body, accent=TEAL):
    rect(slide, x, y, w, h, PANEL, LINE, True)
    rect(slide, x+.18, y+.18, .06, h-.36, accent)
    add_text(slide, label, x+.4, y+.27, w-.65, .2, 8, accent, True, MONO)
    add_text(slide, heading, x+.4, y+.58, w-.65, .38, 17, INK, True)
    add_text(slide, body, x+.4, y+1.05, w-.65, h-1.25, 11.5, MUTED)

def code_box(slide, x, y, w, h, code):
    rect(slide, x, y, w, h, '08151D', LINE, True)
    add_text(slide, code, x+.25, y+.22, w-.5, h-.4, 12, 'C9F3EB', False, MONO)

def make_poster():
    im = Image.new('RGB', (2480, 3508), '#071924')
    draw = ImageDraw.Draw(im)
    font_dir = Path('C:/Windows/Fonts')
    reg = str(font_dir/'malgun.ttf'); bold = str(font_dir/'malgunbd.ttf'); mono = str(font_dir/'consola.ttf')
    def f(path, size): return ImageFont.truetype(path, size)
    # Grid and signal path make the poster readable from a distance.
    for x in range(0, 2480, 80): draw.line((x, 0, x, 3508), fill='#0b2b38', width=1)
    for y in range(0, 3508, 80): draw.line((0, y, 2480, y), fill='#0b2b38', width=1)
    draw.rectangle((0, 0, 30, 3508), fill='#58e0cb')
    draw.text((110, 100), 'SECURE MESSAGE  /  DISCRETE MATHEMATICS LAB', font=f(mono, 28), fill='#58e0cb')
    draw.text((110, 170), '수학으로 설계한\n안전한 메시지 전송', font=f(bold, 106), fill='#eaf5f3', spacing=8)
    draw.text((115, 450), '허프만 트리 · AES-GCM · XOR 학습 모드 · 다익스트라 알고리즘을\n하나의 웹앱으로 구현한 클라이언트 중심 보안 통신 프로젝트', font=f(reg, 38), fill='#a6c3c7', spacing=10)
    # signal path
    pts = [(180, 700), (570, 700), (820, 840), (1210, 840), (1470, 700), (2240, 700)]
    draw.line(pts, fill='#58e0cb', width=12)
    for i, (x, y) in enumerate(pts):
        draw.ellipse((x-34,y-34,x+34,y+34), fill='#f2c85b' if i in (0,len(pts)-1) else '#58e0cb')
    for x, label in [(140,'원본'),(450,'허프만\n압축'),(980,'PBKDF2 +\nAES-GCM'),(1650,'.enc\n패키지'),(2110,'복원')]:
        draw.text((x, 910), label, font=f(bold, 30), fill='#eaf5f3', align='center')
    def poster_card(x,y,w,h,kicker,head,body,color='#58e0cb'):
        draw.rounded_rectangle((x,y,x+w,y+h), radius=28, fill='#103442', outline='#2b5964', width=3)
        draw.rectangle((x+30,y+30,x+40,y+h-30), fill=color)
        draw.text((x+65,y+35), kicker, font=f(mono,22), fill=color)
        draw.text((x+65,y+78), head, font=f(bold,42), fill='#eaf5f3')
        draw.multiline_text((x+65,y+145), body, font=f(reg,27), fill='#a6c3c7', spacing=10)
    draw.rounded_rectangle((110,1120,2370,1450), radius=28, fill='#0c2734', outline='#2b5964', width=3)
    draw.rectangle((145,1155,155,1415), fill='#58e0cb')
    draw.text((190,1160),'02 / WHAT IS THIS PROGRAM?', font=f(mono,25), fill='#58e0cb')
    draw.text((190,1210),'텍스트·파일을 압축하고 암호화하는 웹앱', font=f(bold,50), fill='#eaf5f3')
    draw.multiline_text((190,1285),'입력 → 허프만 압축 → PBKDF2 + AES-256-GCM → .enc 패키지 → 인증 복호화\n원본과 비밀번호는 브라우저에서 처리하고, 공유할 때는 암호문만 업로드한다.', font=f(reg,29), fill='#a6c3c7', spacing=9)

    draw.text((110,1515),'학습한 개념은 실생활에서 어떻게 쓰일까?', font=f(bold,50), fill='#eaf5f3')
    poster_card(110,1600,1100,385,'01 / TREE + DATA','허프만 코딩','파일 형식 ZIP·PNG 등의 무손실 압축처럼\n자주 나타나는 데이터에 짧은 코드를 부여해\n저장 공간과 전송량을 줄인다.', '#58e0cb')
    poster_card(1270,1600,1100,385,'02 / SECURITY','AES-GCM 인증 암호','메신저·클라우드·금융 서비스처럼\n내용을 숨기고, 암호문이 바뀌면\n복호화를 거부해 변조를 확인한다.', '#f2c85b')
    poster_card(110,2035,1100,385,'03 / BOOLEAN ALGEBRA','XOR 학습 모드','A ⊕ K ⊕ K = A 성질은\n암호 알고리즘의 비트 연산 원리를 익히는\n교육용 비교 실험에 사용한다.', '#ff8977')
    poster_card(1270,2035,1100,385,'04 / GRAPH','다익스트라 알고리즘','내비게이션의 최단 시간 길찾기,\n통신망의 비용 계산처럼 가중치 합이 작은\n경로를 찾는 문제에 적용된다.', '#58e0cb')

    draw.text((110,2500),'프로그램 작동 구조', font=f(bold,50), fill='#eaf5f3')
    items = [('입력','텍스트 또는 25MB 이하 파일'),('압축','빈도표 → 허프만 트리'),('키 유도','PBKDF2-SHA-256 250,000회'),('암호화','AES-256-GCM + 인증 태그'),('공유','암호문 .enc만 업로드')]
    x=110
    for i,(a,b) in enumerate(items):
        draw.rounded_rectangle((x,2590,x+430,2840), radius=20, fill='#0c2734', outline='#2b5964', width=2)
        draw.text((x+28,2620), f'{i+1:02d}', font=f(mono,25), fill='#58e0cb')
        draw.text((x+28,2665), a, font=f(bold,31), fill='#eaf5f3')
        draw.multiline_text((x+28,2720), b, font=f(reg,21), fill='#a6c3c7', spacing=6)
        if i < 4: draw.polygon([(x+438,2702),(x+475,2724),(x+438,2746)], fill='#f2c85b')
        x += 470
    draw.rounded_rectangle((110,2910,2370,3290), radius=28, fill='#0c2734', outline='#2b5964', width=3)
    draw.text((155,2950),'실생활 서비스보다 더 나은 개선 제안', font=f(bold,43), fill='#f2c85b')
    improvements=[('접근 제어','공개 보관함 대신 로그인·수신자별 권한을 적용'),('자동 삭제','만료 시 DB 기록뿐 아니라 저장 파일도 함께 삭제'),('키 보호','강한 비밀번호 안내와 키 공유 경고를 제공')]
    for i,(head,body) in enumerate(improvements):
        x=155+i*720
        draw.ellipse((x,3033,x+46,3079), fill='#58e0cb')
        draw.text((x+14,3040), str(i+1), font=f(mono,21), fill='#071924')
        draw.text((x+65,3032), head, font=f(bold,29), fill='#eaf5f3')
        draw.multiline_text((x+65,3080), body, font=f(reg,22), fill='#a6c3c7', spacing=7)
    draw.text((110,3350),'핵심 질문  |  “수학적 구조는 어떻게 안전한 정보 전달을 가능하게 하는가?”', font=f(bold,29), fill='#58e0cb')
    draw.text((110,3420),'SECURE MESSAGE · HIGH SCHOOL DISCRETE MATHEMATICS PROJECT', font=f(mono,21), fill='#a6c3c7')
    path = OUT/'secure-message-poster.png'; im.save(path, quality=95); return path

def slide_deck():
    prs = Presentation(); prs.slide_width = Inches(W); prs.slide_height = Inches(H)
    blank = prs.slide_layouts[6]
    # 1 cover
    s=prs.slides.add_slide(blank); bg(s, 'DISCRETE MATHEMATICS PROJECT / 2026')
    add_text(s,'SECURE',.72,1.13,5.3,.65,45,INK,True,MONO)
    add_text(s,'MESSAGE',.72,1.83,7.2,.72,48,TEAL,True,MONO)
    add_text(s,'수학으로 설계한 안전한 메시지 전송 웹앱',.76,2.75,9,.5,23,INK,True)
    add_text(s,'허프만 트리 · AES-GCM · XOR 학습 모드 · 다익스트라 알고리즘',.78,3.38,8.8,.3,13,MUTED)
    # right network
    nodes=[(10.1,1.4,'입력',GOLD),(11.5,2.5,'압축',TEAL),(9.7,3.7,'AES',CORAL),(11.4,4.8,'.enc',TEAL),(9.9,5.9,'복원',GOLD)]
    for a,b in zip(nodes,nodes[1:]): line(s,a[0]+.35,a[1]+.35,b[0]+.35,b[1]+.35,TEAL,2)
    for x,y,label,c in nodes: circle(s,x,y,.7,c,label,INK,10)
    rect(s,.76,5.35,5.55,.95,PANEL,LINE,True); add_text(s,'핵심 원칙',1.05,5.6,1,.2,9,TEAL,True,MONO); add_text(s,'원본 파일과 비밀 키는 브라우저에서만 처리',2.1,5.57,3.8,.28,14,INK,True)
    footer(s,1)
    # 2 agenda
    s=prs.slides.add_slide(blank); bg(s); title(s,'ROADMAP','발표 흐름','문제 정의부터 수학적 모델, 구현과 한계까지 설명합니다.')
    agenda=[('01','문제와 목표','왜 보안 전송을 수학적으로 모델링하는가'),('02','전체 구조','입력부터 복호화·공유까지의 데이터 흐름'),('03','핵심 개념','허프만 / AES-GCM / XOR 학습 / 다익스트라'),('04','구현과 시연','웹앱 기능과 검증 과정'),('05','한계와 확장','공개 보관함의 한계 및 개선 방향')]
    for i,(n,h,b) in enumerate(agenda):
        x=.75+(i%2)*6.05; y=2.5+(i//2)*1.25
        if i==4: x=3.78
        rect(s,x,y,5.45,.92,PANEL,LINE,True); add_text(s,n,x+.25,y+.22,.45,.2,11,TEAL,True,MONO); add_text(s,h,x+.9,y+.17,1.5,.25,15,INK,True); add_text(s,b,x+.9,y+.49,4.2,.2,10,MUTED)
    footer(s,2)
    # 3 problem
    s=prs.slides.add_slide(blank); bg(s); title(s,'01 / PROBLEM','메시지를 “안전하게” 전송하려면?','기밀성만이 아니라 용량, 무결성, 전달 경로까지 함께 고려해야 합니다.')
    card(s,.7,2.55,3.8,2.45,'CHALLENGE A','용량','파일 크기가 크면 저장·전송 효율이 떨어진다.\n→ 빈도 기반 무손실 압축',TEAL)
    card(s,4.76,2.55,3.8,2.45,'CHALLENGE B','기밀성·무결성','전송 중 내용을 읽거나 바꾸지 못해야 한다.\n→ PBKDF2 + AES-GCM',GOLD)
    card(s,8.82,2.55,3.8,2.45,'CHALLENGE C','경로','같은 데이터라도 비용 기준에 따라\n좋은 전송 경로가 달라진다.',CORAL)
    code_box(s,2.25,5.5,8.85,.65,'목표:  이산수학의 트리 · 부울대수 · 그래프를 실제 동작하는 웹앱에 연결한다.')
    footer(s,3)
    # 4 architecture
    s=prs.slides.add_slide(blank); bg(s); title(s,'02 / ARCHITECTURE','데이터는 이렇게 이동합니다','브라우저에서 암호화한 .enc 패키지만 선택적으로 공유 서버로 전송합니다.')
    stages=[('원본 입력','텍스트 / 파일','01',GOLD),('허프만 압축','빈도표와 트리','02',TEAL),('키 유도·암호화','PBKDF2 + AES-GCM','03',CORAL),('.enc 패키지','헤더 + 인증 암호문','04',TEAL),('복호화·복원','인증 검증','05',GOLD)]
    for i,(h,b,n,c) in enumerate(stages):
        x=.65+i*2.52; rect(s,x,2.8,2.05,1.55,PANEL,LINE,True); circle(s,x+.16,2.98,.36,c,n,NAVY,8); add_text(s,h,x+.2,3.48,1.65,.28,13,INK,True); add_text(s,b,x+.2,3.86,1.65,.3,9,MUTED)
        if i<4: line(s,x+2.05,3.58,x+2.48,3.58,GOLD,2)
    rect(s,.82,5.05,5.35,1.1,'0C2734',LINE,True); add_text(s,'로컬 영역',1.08,5.3,1.2,.2,10,TEAL,True,MONO); add_text(s,'원본 · 비밀 키 · 압축 · 복호화',2.2,5.28,3.5,.25,15,INK,True)
    rect(s,6.55,5.05,5.85,1.1,'0C2734',LINE,True); add_text(s,'공유 선택 시',6.82,5.3,1.4,.2,10,GOLD,True,MONO); add_text(s,'Supabase에는 .enc 암호문과 최소 메타데이터만',8.32,5.28,3.65,.3,13,INK,True)
    footer(s,4)
    # 5 huffman
    s=prs.slides.add_slide(blank); bg(s); title(s,'03 / TREE','허프만 코딩: 빈도가 코드 길이를 결정한다','최소 힙으로 빈도가 가장 작은 두 노드를 반복 결합해 이진 트리를 만듭니다.')
    # tree sample
    coords={'R':(2.9,2.55),'A':(1.75,3.45),'B':(4.15,3.45),'C':(1.1,4.55),'D':(2.4,4.55),'E':(3.55,4.55),'F':(4.8,4.55)}
    for a,b,label in [('R','A','0'),('R','B','1'),('A','C','0'),('A','D','1'),('B','E','0'),('B','F','1')]:
        x1,y1=coords[a];x2,y2=coords[b];line(s,x1,y1,x2,y2,TEAL,1.5); add_text(s,label,(x1+x2)/2-.08,(y1+y2)/2-.16,.2,.2,9,GOLD,True,MONO)
    for key,lab,c in [('R','10',TEAL),('A','4',TEAL),('B','6',TEAL),('C','A:1',GOLD),('D','B:3',GOLD),('E','C:2',GOLD),('F','D:4',GOLD)]:
        x,y=coords[key];circle(s,x-.27,y-.27,.54,c,lab,NAVY if c==TEAL else NAVY,8)
    card(s,6.1,2.65,2.85,2.6,'ALGORITHM','구성 과정','1. 바이트 빈도 계산\n2. 최소 힙에 삽입\n3. 최소 두 노드 결합\n4. 루트까지 반복',TEAL)
    card(s,9.3,2.65,2.85,2.6,'RESULT','특징','자주 등장하는 바이트\n→ 짧은 비트열\n\n모든 코드가 다른 코드의\n접두어가 되지 않음',GOLD)
    code_box(s,1.0,5.75,11.45,.52,'예시: A = 00, B = 01, C = 10, D = 11  →  경계 표시 없이도 원문을 복원 가능')
    footer(s,5)
    # 6 crypto modes
    s=prs.slides.add_slide(blank); bg(s); title(s,'03 / CRYPTOGRAPHY','기본은 AES-GCM, XOR은 원리 학습용','실제 패키지는 PBKDF2-SHA-256으로 키를 유도한 뒤 AES-256-GCM으로 인증 암호화합니다.')
    code_box(s,.9,2.55,11.55,.83,'A ⊕ K ⊕ K = A    ·    (A ⊕ K) ⊕ K = A')
    items=[('비밀번호 + 새 salt','PBKDF2-SHA-256  (250,000회)',GOLD),('새 IV + 압축 데이터','AES-256-GCM + 헤더 인증',TEAL),('SME2 암호문','변조 시 복호화 거부',CORAL)]
    for i,(h,b,c) in enumerate(items):
        y=3.75+i*.7; rect(s,2.0,y,9.25,.47,'0C2734',LINE,True); add_text(s,h,2.25,y+.11,2.1,.2,11,c,True); add_text(s,b,5.0,y+.1,4.7,.22,14,INK,False,MONO)
        if i<2:add_text(s,'⊕',1.38,y+.52,.35,.3,18,TEAL,True,MONO)
    add_text(s,'XOR 학습 모드',.9,6.03,1.35,.2,10,TEAL,True,MONO); add_text(s,'A ⊕ K ⊕ K = A 성질을 관찰하기 위한 호환용 SME1 모드이며, 새 파일의 기본값은 아닙니다.',2.4,6.0,9.5,.28,12,MUTED)
    footer(s,6)
    # 7 integrity and package
    s=prs.slides.add_slide(blank); bg(s); title(s,'04 / PACKAGE','SME2 .enc 패키지와 인증 검증','복호화에 필요한 압축 정보, 키 유도 정보, 인증 암호문을 하나의 패키지로 전달합니다.')
    fields=[('SME2','식별자',TEAL),('이름·MIME','원본 정보',GOLD),('빈도표','허프만 트리 재구성',TEAL),('salt · IV','PBKDF2/AES-GCM 값',CORAL),('암호문 + 태그','변조 검출',GOLD)]
    x=.75
    for h,b,c in fields:
        rect(s,x,2.75,2.28,1.25,PANEL,LINE,True); add_text(s,h,x+.18,3.0,1.85,.23,14,c,True,MONO); add_text(s,b,x+.18,3.39,1.85,.25,10,MUTED); x+=2.48
    add_text(s,'복호화 검증 흐름',.75,4.75,2,.25,15,INK,True)
    for i,t in enumerate(['AES-GCM 인증 복호화','허프만 압축 해제','원본 바이트 복원','변조·오류 시 즉시 거부']):
        x=.8+i*3.05; circle(s,x,5.35,.42,TEAL,str(i+1),NAVY,9); add_text(s,t,x+.56,5.42,2.25,.22,11,INK,True)
        if i<3:line(s,x+2.55,5.56,x+3.0,5.56,GOLD,1.5)
    footer(s,7)
    # 8 dijkstra
    s=prs.slides.add_slide(blank); bg(s); title(s,'03 / GRAPH','다익스트라: 가중치 합이 가장 작은 경로 찾기','정점은 서버, 간선은 연결, 가중치는 전송 시간 또는 위험도로 모델링했습니다.')
    pts=[(1.2,4.2,'나'),(3.05,2.8,'서울'),(3.15,5.0,'대전'),(5.35,2.65,'인천'),(5.8,4.9,'부산'),(8.0,3.65,'광주'),(10.5,3.75,'상대')]
    edge_idx=[(0,1,12),(0,2,20),(1,2,11),(1,3,9),(2,3,16),(2,4,13),(3,4,14),(3,5,18),(4,5,10),(4,6,23),(5,6,8)]
    highlight={(0,1),(1,3),(3,5),(5,6)}
    for a,b,wgt in edge_idx:
        x1,y1,_=pts[a];x2,y2,_=pts[b];c=TEAL if (a,b) in highlight else LINE;line(s,x1+.3,y1+.3,x2+.3,y2+.3,c,2.6 if c==TEAL else 1.1);add_text(s,str(wgt),(x1+x2)/2+.12,(y1+y2)/2-.12,.3,.18,8,GOLD,True,MONO)
    for x,y,label in pts:circle(s,x,y,.6,GOLD if label in ('나','상대') else TEAL,label,NAVY,10)
    card(s,.85,5.65,3.6,.85,'MODEL','시간 가중치','선 위 숫자는 ms. 합이 작은 경로를 선택.',TEAL)
    card(s,4.85,5.65,3.6,.85,'MODEL','위험도 가중치','같은 그래프도 비용 정의에 따라 결과가 달라짐.',GOLD)
    card(s,8.85,5.65,3.6,.85,'VISUALIZATION','시연','선택 경로와 탐색 기록, 패킷 이동을 표시.',CORAL)
    footer(s,8)
    # 9 live features
    s=prs.slides.add_slide(blank); bg(s); title(s,'05 / IMPLEMENTATION','웹앱에서 직접 확인할 수 있는 기능','입력-처리-시각화-공유가 하나의 흐름으로 연결됩니다.')
    features=[('암호화','PBKDF2 + AES-GCM 기본, 25MB 제한, .enc 다운로드',TEAL),('복호화','동일 키 입력, 인증 검증 후 원본 복원',GOLD),('전송 경로','시간/위험도 선택, 가중치 수정과 비교',CORAL),('보안 분석','XOR 학습용 바이트 빈도 분포 비교',TEAL),('공유 보관함','Supabase에 암호문만 업로드·수신',GOLD)]
    for i,(h,b,c) in enumerate(features):
        x=.75+(i%2)*6.15;y=2.55+(i//2)*1.2
        if i==4:x=3.83
        card(s,x,y,5.65,.92,f'{i+1:02d} / FEATURE',h,b,c)
    footer(s,9)
    # 10 demo
    s=prs.slides.add_slide(blank); bg(s); title(s,'06 / DEMO','시연은 90초 안에 이렇게 진행합니다','짧은 문장 하나로 압축, 암호화, 공유, 복호화까지 보여 줄 수 있습니다.')
    demo=[('1','입력','“이산수학으로 보낸 비밀 메시지”와 비밀번호 입력'),('2','암호화','트리 생성, PBKDF2 키 유도, AES-GCM 인증 암호화'),('3','다운로드/업로드','SME2 .enc 파일 생성 또는 Supabase 보관함 업로드'),('4','경로 계산','시간 또는 위험도 기준으로 다익스트라 실행'),('5','복호화','같은 키로 복원, 인증 검증 통과 확인')]
    for i,(n,h,b) in enumerate(demo):
        y=2.45+i*.78;circle(s,.95,y,.42,TEAL,n,NAVY,9);add_text(s,h,1.6,y+.04,1.15,.22,14,INK,True);add_text(s,b,2.8,y+.06,8.7,.25,11,MUTED);line(s,1.16,y+.43,1.16,y+.73,LINE,1)
    footer(s,10)
    # 11 security
    s=prs.slides.add_slide(blank); bg(s); title(s,'07 / SECURITY','보호되는 것과 남아 있는 한계','검증된 인증 암호를 적용해도 공개 공유 환경의 접근 제어 문제는 별도로 남습니다.')
    card(s,.75,2.65,3.75,2.6,'WHAT IS PROTECTED','보호되는 것','원본 파일과 비밀 키는\n브라우저 내부에서 처리\n\n서버에는 .enc 암호문과\n최소 메타데이터만 저장',TEAL)
    card(s,4.8,2.65,3.75,2.6,'DEFAULT SECURITY','기본 암호화','새 salt + PBKDF2 250,000회\n새 IV + AES-256-GCM\n헤더와 암호문의 변조를 검출',CORAL)
    card(s,8.85,2.65,3.75,2.6,'REMAINING LIMITS','공유 환경의 한계','공개 보관함 정책에서는\n접근 제어가 부족할 수 있음\n권한 관리·삭제·감사 로그 필요',GOLD)
    code_box(s,1.35,5.8,10.65,.58,'결론: “암호화했다”가 아니라 “어떤 위협을 막고, 무엇이 남는가”를 함께 설명한다.')
    footer(s,11)
    # 12 conclusion
    s=prs.slides.add_slide(blank); bg(s,'CONCLUSION / SECURE MESSAGE')
    add_text(s,'TREE',.75,1.15,2.2,.45,27,TEAL,True,MONO); add_text(s,'+',3.02,1.13,.4,.45,27,GOLD,True,MONO);add_text(s,'AES-GCM',3.55,1.15,2.0,.45,27,GOLD,True,MONO);add_text(s,'+',5.75,1.13,.4,.45,27,CORAL,True,MONO);add_text(s,'GRAPH',6.3,1.15,2.8,.45,27,CORAL,True,MONO)
    add_text(s,'수학의 구조를\n“작동하는 보안 통신”으로',.75,2.05,9.5,1.45,38,INK,True)
    add_text(s,'허프만 트리는 효율을, AES-GCM은 기밀성·무결성을, XOR은 부울대수 원리를, 다익스트라는 전달의 최적화를 설명합니다.',.8,3.85,11.4,.35,14,MUTED)
    rect(s,.78,5.0,11.72,1.0,PANEL,LINE,True); add_text(s,'발표 핵심 한 문장',1.05,5.25,1.7,.2,10,TEAL,True,MONO);add_text(s,'“데이터의 표현·변환·이동은 모두 이산수학적 구조로 모델링할 수 있다.”',2.85,5.2,8.8,.32,16,INK,True)
    add_text(s,'THANK YOU',.8,6.4,3,.25,15,TEAL,True,MONO);footer(s,12)
    output=OUT/'secure-message-presentation.pptx';prs.save(output);return output

if __name__ == '__main__':
    poster=make_poster(); deck=slide_deck(); print(poster); print(deck)
