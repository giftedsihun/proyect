# -*- coding: utf-8 -*-
"""Apply the final technical corrections and add a live-app slide."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


SOURCE = Path("secure-message-presentation-final-v3.pptx")
OUTPUT = Path("secure-message-presentation-final-v4.pptx")
APP_URL = "https://eclectic-sunshine-d10d86.netlify.app/"


def replace_text(slide, old, new):
    """Replace a complete text box while preserving its existing styling."""
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False) and shape.text == old:
            paragraph = shape.text_frame.paragraphs[0]
            runs = list(paragraph.runs)
            if runs:
                runs[0].text = new
                for run in runs[1:]:
                    run._r.getparent().remove(run._r)
            else:
                paragraph.text = new
            return
    raise ValueError(f"Text not found: {old!r}")


def add_execution_slide(presentation):
    """Add a minimal final slide whose large label opens the deployed app."""
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(15, 25, 35)

    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.72), Inches(1.4), Inches(0.12), Inches(4.7))
    accent.fill.solid()
    accent.fill.fore_color.rgb = RGBColor(59, 199, 179)
    accent.line.fill.background()

    title = slide.shapes.add_textbox(Inches(1.2), Inches(2.45), Inches(10.8), Inches(1.15))
    frame = title.text_frame
    frame.clear()
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    run = paragraph.add_run()
    run.text = "실행"
    run.font.name = "Pretendard"
    run.font.size = Pt(52)
    run.font.bold = True
    run.font.color.rgb = RGBColor(245, 248, 250)
    run.hyperlink.address = APP_URL

    url_box = slide.shapes.add_textbox(Inches(1.2), Inches(3.8), Inches(10.8), Inches(0.45))
    frame = url_box.text_frame
    frame.clear()
    paragraph = frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    run = paragraph.add_run()
    run.text = APP_URL
    run.font.name = "Pretendard"
    run.font.size = Pt(15)
    run.font.color.rgb = RGBColor(159, 221, 211)
    run.hyperlink.address = APP_URL

    note = slide.shapes.add_textbox(Inches(1.2), Inches(5.2), Inches(10.8), Inches(0.4))
    frame = note.text_frame
    frame.clear()
    paragraph = frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    run = paragraph.add_run()
    run.text = "'실행'을 클릭하면 Secure Message 웹앱이 열립니다."
    run.font.name = "Pretendard"
    run.font.size = Pt(14)
    run.font.color.rgb = RGBColor(180, 190, 200)


def main():
    presentation = Presentation(SOURCE)
    slides = presentation.slides

    # Slide 1: title matches the app's text and file handling scope.
    replace_text(slides[0], "수학으로 설계한 안전한 메시지 전송 웹앱", "수학으로 설계한 안전한 메시지·파일 전송 웹앱")

    # Slide 2: make real-world claims precise rather than universal.
    replace_text(slides[1], "ZIP·PNG의 DEFLATE는 허프만 코딩을 활용", "ZIP·PNG의 DEFLATE는 LZ77과 허프만 코딩을 함께 활용")
    replace_text(slides[1], "HTTPS·파일 보호에는 AES-GCM 같은 인증 암호 사용", "HTTPS·파일 보호에는 AES-GCM 같은 인증 암호 방식이 사용될 수 있음")

    # Slide 5: align the visible Huffman example with variable code lengths.
    replace_text(slides[4], "표시된 트리 예시: A = 00, B = 01, C = 10, D = 11 → 경계 표시 없이 복원 가능", "예시: 빈도 A=4, B=3, C=2, D=1이면 A=0, B=10, C=110, D=111 → 경계 없이 복원 가능")

    # Slides 6-11: fix section sequence, clarify XOR, and reduce the impression
    # that Dijkstra is a current feature.
    replace_text(slides[5], "03 / CRYPTOGRAPHY", "04 / CRYPTOGRAPHY")
    replace_text(slides[5], "A ⊕ K ⊕ K = A 성질을 관찰하기 위한 호환용 SME1 모드이며, 새 파일의 기본값은 아닙니다.", "A ⊕ K ⊕ K = A 성질을 관찰하는 학습·비교용 방식입니다. 반복 XOR은 안전한 실사용 암호가 아니며, 새 파일의 기본값도 아닙니다.")
    replace_text(slides[6], "04 / PACKAGE", "05 / PACKAGE")
    replace_text(slides[7], "05 / IMPLEMENTATION", "06 / IMPLEMENTATION")
    replace_text(slides[8], "06 / DEMO", "07 / DEMO")
    replace_text(slides[8], "시연은 90초 안에 이렇게 진행합니다", "시연은 짧은 텍스트로 이렇게 진행합니다")
    replace_text(slides[9], "07 / SECURITY", "08 / SECURITY")
    replace_text(slides[9], "원본 파일과 비밀 키는\n브라우저 내부에서 처리\n\n서버에는 .enc 암호문과\n최소 메타데이터만 저장", "원본 파일과 비밀 키는\n브라우저 내부에서 처리\n\n서버에는 .enc 암호문과\n최소 메타데이터만 저장\n\n파일명·크기·시각 등 메타데이터는\n별도 보호가 필요")
    replace_text(slides[10], "08 / FUTURE IMPROVEMENT", "09 / FUTURE IMPROVEMENT")
    replace_text(slides[10], "향후 개선 - 다익스트라로 전송 경로 최적화", "향후 확장 - 다익스트라 전송 경로 모델")
    replace_text(slides[10], "현재 웹앱에는 미구현이며, 서버·연결·비용을 그래프로 모델링해 확장할 수 있습니다.", "현재 웹앱에는 미구현입니다. 인증·권한·자동 삭제를 우선 보완한 뒤, 서버·연결·비용을 그래프로 모델링해 확장할 수 있습니다.")
    replace_text(slides[10], "시간 또는 위험도 기준으로 최저 비용 경로를 탐색.", "이산수학 확장 예시: 비용 기준에 따른 최저 경로 탐색.")

    add_execution_slide(presentation)
    presentation.save(OUTPUT)
    print(OUTPUT.resolve())


if __name__ == "__main__":
    main()
